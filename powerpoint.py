from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import pandas as pd

# Colors
DARK_BLUE = RGBColor(26, 35, 126)
MED_BLUE  = RGBColor(21, 101, 192)
WHITE     = RGBColor(255, 255, 255)
GREEN     = RGBColor(46, 125, 50)
RED       = RGBColor(198, 40, 40)
GRAY      = RGBColor(55, 71, 79)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

blank = prs.slide_layouts[6]  # blank layout

def add_rect(slide, l, t, w, h, color):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color

# ── Slide 1: Title ──
slide1 = prs.slides.add_slide(blank)
add_rect(slide1, 0, 0, 13.33, 7.5, DARK_BLUE)
add_rect(slide1, 0, 2.8, 13.33, 0.08, MED_BLUE)
add_text(slide1, "NFHS-5 India Healthcare Analysis",
         1, 1.5, 11, 1.5, size=36, bold=True, align=PP_ALIGN.CENTER)
add_text(slide1, "National Family Health Survey 2019–2021 | State-wise Study",
         1, 3.2, 11, 1, size=18, align=PP_ALIGN.CENTER)
add_text(slide1, "Prepared by: Lavish Pratap Singh",
         1, 4.5, 11, 0.8, size=16, align=PP_ALIGN.CENTER)
add_text(slide1, "Source: data.gov.in | Department of Health & Family Welfare",
         1, 5.5, 11, 0.8, size=13, align=PP_ALIGN.CENTER)

# ── Slide 2: Agenda ──
slide2 = prs.slides.add_slide(blank)
add_rect(slide2, 0, 0, 13.33, 1.2, DARK_BLUE)
add_text(slide2, "Agenda", 0.3, 0.2, 10, 0.8, size=28, bold=True, align=PP_ALIGN.LEFT)

points = [
    "01  Project Overview & Objectives",
    "02  Data Source — NFHS-5 (2019-21)",
    "03  Female Literacy Analysis",
    "04  Health Insurance Coverage",
    "05  Sanitation Access",
    "06  Key Findings & Recommendations",
]
for i, point in enumerate(points):
    bg = MED_BLUE if i % 2 == 0 else GRAY
    add_rect(slide2, 0.3, 1.4 + i*0.9, 12.5, 0.8, bg)
    add_text(slide2, point, 0.6, 1.5 + i*0.9, 12, 0.7, size=16, bold=False)

# ── Slide 3: Overview ──
slide3 = prs.slides.add_slide(blank)
add_rect(slide3, 0, 0, 13.33, 1.2, DARK_BLUE)
add_text(slide3, "Project Overview", 0.3, 0.2, 10, 0.8, size=28, bold=True)

# KPI boxes
kpis = [("36", "States\nAnalyzed"), ("5", "Key\nIndicators"),
        ("2019-21", "Survey\nPeriod"), ("136", "Data\nColumns")]
kpi_colors = [MED_BLUE, GREEN, GRAY, RED]
for i, ((val, label), color) in enumerate(zip(kpis, kpi_colors)):
    x = 0.5 + i * 3.1
    add_rect(slide3, x, 1.5, 2.8, 2.2, color)
    add_text(slide3, val, x, 1.7, 2.8, 1.0, size=32, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide3, label, x, 2.8, 2.8, 0.8, size=14, align=PP_ALIGN.CENTER)

add_text(slide3, "This study analyzes NFHS-5 data to understand healthcare disparities\nacross Indian states covering literacy, sanitation, sex ratio and insurance.",
         0.5, 4.2, 12, 1.5, size=15, color=GRAY, align=PP_ALIGN.CENTER)

# ── Slide 4: Female Literacy Chart ──
slide4 = prs.slides.add_slide(blank)
add_rect(slide4, 0, 0, 13.33, 1.2, DARK_BLUE)
add_text(slide4, "Female Literacy — State Analysis", 0.3, 0.2, 12, 0.8, size=28, bold=True)
try:
    pic = slide4.shapes.add_picture('chart_literacy.png', Inches(1), Inches(1.4), Inches(11), Inches(5.5))
except:
    add_text(slide4, "chart_literacy.png not found", 1, 3, 11, 1, size=16, color=RED, align=PP_ALIGN.CENTER)

# ── Slide 5: Health Insurance Chart ──
slide5 = prs.slides.add_slide(blank)
add_rect(slide5, 0, 0, 13.33, 1.2, DARK_BLUE)
add_text(slide5, "Health Insurance Coverage", 0.3, 0.2, 12, 0.8, size=28, bold=True)
try:
    pic = slide5.shapes.add_picture('chart_insurance.png', Inches(1), Inches(1.4), Inches(11), Inches(5.5))
except:
    add_text(slide5, "chart_insurance.png not found", 1, 3, 11, 1, size=16, color=RED, align=PP_ALIGN.CENTER)

# ── Slide 6: Sanitation Chart ──
slide6 = prs.slides.add_slide(blank)
add_rect(slide6, 0, 0, 13.33, 1.2, DARK_BLUE)
add_text(slide6, "Sanitation Access — Lowest States", 0.3, 0.2, 12, 0.8, size=28, bold=True)
try:
    pic = slide6.shapes.add_picture('chart_sanitation.png', Inches(1), Inches(1.4), Inches(11), Inches(5.5))
except:
    add_text(slide6, "chart_sanitation.png not found", 1, 3, 11, 1, size=16, color=RED, align=PP_ALIGN.CENTER)

# ── Slide 7: Key Findings ──
slide7 = prs.slides.add_slide(blank)
add_rect(slide7, 0, 0, 13.33, 1.2, DARK_BLUE)
add_text(slide7, "Key Findings", 0.3, 0.2, 12, 0.8, size=28, bold=True)

findings = [
    ("🔴", "Bihar female literacy only 57% — needs urgent attention", RED),
    ("🔴", "Sanitation access in Bihar just 49% — half households uncovered", RED),
    ("🟡", "Health insurance national avg only 41% — huge gap", RGBColor(245, 127, 23)),
    ("🟢", "Kerala & Lakshadweep lead in female literacy (95%+)", GREEN),
    ("🟢", "Goa & Andaman top in sanitation access (87%+)", GREEN),
]
for i, (icon, text, color) in enumerate(findings):
    add_rect(slide7, 0.4, 1.4 + i*1.0, 12.3, 0.85, RGBColor(236, 239, 241))
    add_text(slide7, icon + "  " + text, 0.7, 1.5 + i*1.0, 11.5, 0.75,
             size=15, color=color, bold=False)

# ── Slide 8: Recommendations ──
slide8 = prs.slides.add_slide(blank)
add_rect(slide8, 0, 0, 13.33, 1.2, DARK_BLUE)
add_text(slide8, "Recommendations", 0.3, 0.2, 12, 0.8, size=28, bold=True)

recs = [
    ("HIGH", "Launch girl education campaigns in Bihar, Rajasthan, UP", RED),
    ("HIGH", "Accelerate Swachh Bharat Mission in low-sanitation states", RED),
    ("MED",  "Expand Ayushman Bharat PM-JAY coverage nationally", RGBColor(245, 127, 23)),
    ("MED",  "Strengthen Beti Bachao Beti Padhao in low sex-ratio states", RGBColor(245, 127, 23)),
    ("LOW",  "Adopt Kerala healthcare model in lagging states", GREEN),
]
for i, (priority, text, color) in enumerate(recs):
    add_rect(slide8, 0.4, 1.4 + i*1.0, 1.5, 0.8, color)
    add_text(slide8, priority, 0.4, 1.45 + i*1.0, 1.5, 0.7,
             size=12, bold=True, align=PP_ALIGN.CENTER)
    add_rect(slide8, 2.1, 1.4 + i*1.0, 10.8, 0.8, RGBColor(236, 239, 241))
    add_text(slide8, text, 2.3, 1.48 + i*1.0, 10.5, 0.7,
             size=14, color=GRAY)

# ── Slide 9: Thank You ──
slide9 = prs.slides.add_slide(blank)
add_rect(slide9, 0, 0, 13.33, 7.5, DARK_BLUE)
add_rect(slide9, 0, 3.2, 13.33, 0.06, MED_BLUE)
add_text(slide9, "Thank You", 0, 1.5, 13.33, 1.5,
         size=48, bold=True, align=PP_ALIGN.CENTER)
add_text(slide9, "NFHS-5 Healthcare Analysis | data.gov.in",
         0, 3.5, 13.33, 1, size=18, align=PP_ALIGN.CENTER)
add_text(slide9, "Lavish Pratap Singh", 0, 5, 13.33, 0.8,
         size=16, align=PP_ALIGN.CENTER)

prs.save('Healthcare_Presentation.pptx')
print("PowerPoint ready! Healthcare_Presentation.pptx are ready!")